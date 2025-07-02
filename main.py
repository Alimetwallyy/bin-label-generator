import streamlit as st
import pandas as pd
import io
import plotly.graph_objects as go
import seaborn as sns
import string
import re
from openpyxl import Workbook
from openpyxl.styles import PatternFill, Alignment, Font, Border, Side
from openpyxl.utils.dataframe import dataframe_to_rows
from openpyxl.utils import get_column_letter
# Added for Tab 4
from pptx import Presentation
from pptx.util import Inches, Cm, Pt
from pptx.enum.shapes import MSO_SHAPE


# Add "Created By Alimomet" in top left
st.markdown("""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Roboto:wght@700&display=swap');
    .created-by {
        position: absolute;
        top: 10px;
        left: 20px;
        font-family: 'Roboto', Arial, Helvetica, sans-serif;
        font-size: 14px;
        font-weight: bold;
        color: #333333;
        z-index: 1000;
    }
    </style>
    <div class="created-by">Created By Alimomet</div>
""", unsafe_allow_html=True)

def generate_bin_labels_table(group_name, bay_ids, shelves, bins_per_shelf):
    data = []
    for bay in bay_ids:
        try:
            base_label = bay.replace("BAY-", "")
            base_number = int(base_label[-3:])
            aisle_match = re.search(r'\d{3}', base_label)
            aisle = aisle_match.group(0) if aisle_match else ""

            max_bins = max(bins_per_shelf.get(shelf, 0) for shelf in shelves) if shelves else 1

            for i in range(max_bins):
                row = {
                    'BAY TYPE': group_name,
                    'AISLE': aisle,
                    'BAY ID': bay
                }
                for shelf in shelves:
                    shelf_bin_count = bins_per_shelf.get(shelf, 0)
                    if i < shelf_bin_count:
                        bin_label = base_label[:-4] + shelf + f"{base_number + i:03d}"
                        row[shelf] = bin_label
                    else:
                        row[shelf] = None
                data.append(row)
        except Exception as e:
            st.error(f"Error processing bay ID '{bay}': {str(e)}")
    return pd.DataFrame(data)

def plot_bin_diagram(bay_id, shelves, bins_per_shelf, base_number):
    try:
        fig = go.Figure()
        colors = sns.color_palette("colorblind", len(shelves) if shelves else 1).as_hex()
        shelf_colors = {shelf: colors[i % len(colors)] for i, shelf in enumerate(shelves)} if shelves else {}

        for col_idx, shelf in enumerate(shelves):
            shelf_bins = bins_per_shelf.get(shelf, 0)
            for i in range(shelf_bins):
                bin_label = bay_id.replace("BAY-", "")[:-4] + shelf + f"{base_number + i:03d}"
                x0, x1 = col_idx - 0.4, col_idx + 0.4
                y0, y1 = -i - 0.4, -i + 0.4
                fig.add_shape(
                    type="rect",
                    x0=x0,
                    x1=x1,
                    y0=y0,
                    y1=y1,
                    fillcolor=shelf_colors.get(shelf, "lightblue"),
                    line=dict(color="black"),
                    label=dict(text=bin_label, textposition="middle center", font=dict(size=10)),
                )
                fig.add_trace(
                    go.Scatter(
                        x=[(x0 + x1) / 2],
                        y=[(y0 + y1) / 2],
                        text=[bin_label],
                        mode="text",
                        hoverinfo="text",
                        showlegend=False,
                    )
                )

        fig.update_layout(
            title=f"Bin Layout for {bay_id}",
            xaxis=dict(
                tickmode="array",
                tickvals=list(range(len(shelves))) if shelves else [0],
                ticktext=shelves if shelves else ["No Shelves"],
                showgrid=False,
                zeroline=False,
            ),
            yaxis=dict(
                showgrid=False,
                zeroline=False,
                autorange="reversed",
            ),
            showlegend=bool(shelves),
            legend_title_text="Shelves",
            width=200 * (len(shelves) if shelves else 1),
            height=100 * (max(bins_per_shelf.values(), default=1) if bins_per_shelf else 1),
            margin=dict(l=20, r=20, t=50, b=20),
        )

        for shelf in shelves:
            fig.add_trace(
                go.Scatter(
                    x=[None],
                    y=[None],
                    mode="markers",
                    name=shelf,
                    marker=dict(size=10, color=shelf_colors.get(shelf, "lightblue")),
                )
            )

        return fig
    except Exception as e:
        st.error(f"Error generating diagram for '{bay_id}': {str(e)}")
        return None

def style_excel(writer, sheet_name, df, shelves):
    try:
        ws = writer.sheets[sheet_name]
        yellow_fill = PatternFill(start_color="FFFF00", end_color="FFFF00", fill_type="solid")
        border = Border(left=Side(style='thin'), right=Side(style='thin'),
                        top=Side(style='thin'), bottom=Side(style='thin'))
        bold_font = Font(bold=True)
        center_align = Alignment(horizontal="center", vertical="center")

        hex_colors = [
            "339900", "9B30FF", "FFFF00", "00FFFF", "CC0000", "F88017",
            "FF00FF", "996600", "00FF00", "FF6565", "9999FE"
        ]
        
        styling_colors = ["FFFFFF"] + hex_colors

        if shelves:
            ws.merge_cells('A1:C1')
            ws['A1'] = "HEX COLOR CODES ->"
            ws['A1'].fill = yellow_fill
            ws['A1'].font = bold_font
            ws['A1'].alignment = center_align
            ws['A1'].border = border
            
            for i, hex_color in enumerate(styling_colors[:len(shelves)]):
                col_letter = get_column_letter(4 + i)
                ws[f"{col_letter}1"] = hex_color
                ws[f"{col_letter}1"].fill = PatternFill(start_color=hex_color, end_color=hex_color, fill_type="solid")
                ws[f"{col_letter}1"].font = bold_font
                ws[f"{col_letter}1"].alignment = center_align
                ws[f"{col_letter}1"].border = border

                ws[f"{col_letter}2"] = shelves[i]
                ws[f"{col_letter}2"].fill = PatternFill(start_color=hex_color, end_color=hex_color, fill_type="solid")
                ws[f"{col_letter}2"].font = bold_font
                ws[f"{col_letter}2"].alignment = center_align
                ws[f"{col_letter}2"].border = border

        header_row = 2 if shelves else 1
        for col in range(1, df.shape[1] + 1):
            cell = ws.cell(row=header_row, column=col)
            cell.font = bold_font
            cell.alignment = center_align
            cell.border = border

        for row in ws.iter_rows(min_row=header_row + 1, max_row=ws.max_row, max_col=ws.max_column):
            for cell in row:
                if cell.value is not None:
                    cell.font = bold_font
                    cell.alignment = center_align
                    cell.border = border
    except Exception as e:
        st.error(f"Error styling Excel sheet '{sheet_name}': {str(e)}")

def check_duplicate_bay_ids(bay_groups):
    errors = []
    all_bay_ids = {}

    for group_idx, group in enumerate(bay_groups):
        group_name = group["name"]
        bay_ids = [bay_id.strip().upper() for bay_id in group["bays"] if bay_id.strip()]

        seen_in_group = set()
        for bay_id in bay_ids:
            if bay_id in seen_in_group:
                errors.append(f"⚠️ Duplicate bay ID '{bay_id}' found in {group_name}.")
            seen_in_group.add(bay_id)

            if bay_id not in all_bay_ids:
                all_bay_ids[bay_id] = [group_name]
            elif group_name not in all_bay_ids[bay_id]:
                all_bay_ids[bay_id].append(group_name)

    for bay_id, groups in all_bay_ids.items():
        if len(groups) > 1:
            errors.append(f"⚠️ Bay ID '{bay_id}' is duplicated across groups: {', '.join(groups)}.")

    return errors

def check_duplicate_bin_ids(bay_groups):
    errors = []
    all_bin_ids = {}

    for group_idx, group in enumerate(bay_groups):
        group_name = group["name"]
        bin_ids = [bin_id.strip().upper() for bin_id in group["bin_ids"] if bin_id.strip()]

        seen_in_group = set()
        for bin_id in bin_ids:
            if bin_id in seen_in_group:
                errors.append(f"⚠️ Duplicate bin ID '{bin_id}' found in {group_name}.")
            seen_in_group.add(bin_id)

            if bin_id not in all_bin_ids:
                all_bin_ids[bin_id] = [group_name]
            elif group_name not in all_bin_ids[bin_id]:
                all_bin_ids[bin_id].append(group_name)

    for bin_id, groups in all_bin_ids.items():
        if len(groups) > 1:
            errors.append(f"⚠️ Bin ID '{bin_id}' is duplicated across groups: {', '.join(groups)}.")

    return errors

def parse_bay_definition(bay_definition):
    try:
        if not bay_definition:
            raise ValueError("Bay Definition cannot be empty.")
        return {"bay_definition": bay_definition}
    except Exception as e:
        return {"error": str(e)}

def generate_elevation_powerpoint(bay_types_data):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6] 
    
    for bay_type in bay_types_data:
        slide = prs.slides.add_slide(blank_slide_layout)
        
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.33), Inches(0.5))
        title_text = title_shape.text_frame
        p = title_text.paragraphs[0]
        p.text = bay_type['name']
        p.font.bold = True
        p.font.size = Pt(24)

        start_x = Inches(2.0)
        start_y = Inches(7.0) 
        scale = Cm(0.2)
        current_y = start_y
        
        for i, shelf_name in enumerate(reversed(bay_type['shelves'])):
            shelf_info = bay_type['shelf_details'][shelf_name]
            num_bins = shelf_info['num_bins']
            bin_h = shelf_info['h'] * scale
            bin_w = shelf_info['w'] * scale
            
            shelf_height = bin_h
            shelf_width = num_bins * bin_w

            slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, start_x, current_y - shelf_height, shelf_width, shelf_height)
            
            for j in range(num_bins):
                bin_x = start_x + (j * bin_w)
                slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bin_x, current_y - shelf_height, bin_w, shelf_height)

            label_box = slide.shapes.add_textbox(start_x - Inches(0.6), current_y - shelf_height, Inches(0.5), shelf_height)
            label_box.text_frame.text = shelf_name
            label_box.text_frame.paragraphs[0].font.size = Pt(10)

            is_first_of_dim = True
            if i > 0:
                prev_shelf_name = list(reversed(bay_type['shelves']))[i-1]
                if bay_type['shelf_details'][prev_shelf_name] == shelf_info:
                    is_first_of_dim = False
            
            if is_first_of_dim:
                dim_x = start_x + shelf_width + Inches(0.2)
                dim_y = current_y - (shelf_height / 2)
                dim_text = f"H: {shelf_info['h']}cm\nW: {shelf_info['w']}cm\nD: {shelf_info['d']}cm"
                dim_box = slide.shapes.add_textbox(dim_x, dim_y - Inches(0.25), Inches(1.5), Inches(0.5))
                tf = dim_box.text_frame
                tf.text = dim_text
                for para in tf.paragraphs:
                    para.font.size = Pt(9)

            current_y -= (shelf_height + Cm(0.2))

    ppt_buffer = io.BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    return ppt_buffer

# --- Main App ---
st.set_page_config(layout="wide")

st.title("Space Launch Quick Tools")
st.markdown("A collection of tools for space launch operations.")

if 'eoa_placement_rule' not in st.session_state:
    st.session_state.eoa_placement_rule = "Odd on Left / Even on Right"

tab1, tab2, tab3, tab4 = st.tabs(["Bin Label Generator", "Bin Bay Mapping", "EOA Generator", "Bay Elevation Generator"])

with tab1:
    st.header("Bin Label Generator 🏷️", divider='rainbow')
    st.markdown("Define bay groups, shelves, and bins per shelf to generate structured bin labels. Bay IDs must be unique (e.g., BAY-001-001-001).")

    bay_groups = []
    duplicate_errors = []
    num_groups = st.number_input("How many bay groups do you want to define?", min_value=1, max_value=50, value=1, key="num_groups_bin_label")

    for group_idx in range(num_groups):
        if f"group_name_{group_idx}" not in st.session_state:
            st.session_state[f"group_name_{group_idx}"] = f"Bay Group {group_idx + 1}"

        def update_group_name(idx=group_idx):
            st.session_state[f"group_name_{idx}"] = st.session_state[f"group_name_input_{idx}"]

        header = st.session_state[f"group_name_{group_idx}"].strip() or f"Bay Group {group_idx + 1}"

        with st.expander(header, expanded=True):
            st.text_input(
                "Group Name",
                value=st.session_state[f"group_name_{group_idx}"],
                key=f"group_name_input_{group_idx}",
                on_change=update_group_name,
                args=(group_idx,)
            )

            bays_input = st.text_area(f"Enter bay IDs (one per line, e.g., BAY-001-001-001)", key=f"bays_{group_idx}")
            shelf_count = st.number_input("How many shelves?", min_value=1, max_value=26, value=3, key=f"shelf_count_{group_idx}")
            shelves = list(string.ascii_uppercase[:shelf_count])
            
            st.divider()

            bins_per_shelf = {}
            st.markdown("**Bins per Shelf**")
            for shelf in shelves:
                count = st.number_input(f"Number of bins in shelf {shelf}", min_value=1, max_value=100, value=5, key=f"bins_{group_idx}_{shelf}")
                bins_per_shelf[shelf] = count

            if bays_input:
                bay_list = [b.strip() for b in bays_input.splitlines() if b.strip()]
                if bay_list:
                    bay_groups.append({
                        "name": st.session_state[f"group_name_{group_idx}"].strip() or f"Bay Group {group_idx + 1}",
                        "bays": bay_list,
                        "shelves": shelves,
                        "bins_per_shelf": bins_per_shelf
                    })

    if st.button("Generate Bin Labels", key="generate_bin_labels_full"):
        with st.spinner("Generating bin labels and diagrams..."):
            total_labels_generated = 0
            total_bays_processed = 0
            output = io.BytesIO()
            try:
                with pd.ExcelWriter(output, engine='openpyxl') as writer:
                    for group in bay_groups:
                        df = generate_bin_labels_table(group["name"], group["bays"], group["shelves"], group["bins_per_shelf"])
                        if not df.empty:
                            shelf_cols = [col for col in df.columns if col not in ['BAY TYPE', 'AISLE', 'BAY ID']]
                            total_labels_generated += df[shelf_cols].count().sum()
                            total_bays_processed += df['BAY ID'].nunique()
                            df.to_excel(writer, index=False, startrow=1, sheet_name=group["name"])
                            style_excel(writer, group["name"], df, group["shelves"])
                output.seek(0)
                
                st.success(f"✅ Success! Generated {total_labels_generated} labels for {total_bays_processed} bays across {len(bay_groups)} groups.")
                st.download_button(
                    label="📥 Download Excel File",
                    data=output,
                    file_name="bin_labels.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                )

                st.subheader("🖼️ Interactive Bin Layout Diagrams")
                st.caption("Click on a bay to expand its visual layout.")
                for group in bay_groups:
                    for bay_id in group['bays']:
                        with st.expander(f"View Diagram for **{bay_id}**"):
                            try:
                                base_label = bay_id.replace("BAY-", "")
                                base_number = int(base_label[-3:])
                                fig = plot_bin_diagram(bay_id, group['shelves'], group['bins_per_shelf'], base_number)
                                if fig:
                                    st.plotly_chart(fig, use_container_width=True)
                            except Exception as e:
                                st.error(f"Could not generate diagram for {bay_id}: {e}")
            except Exception as e:
                st.error(f"Error generating output: {str(e)}")

with tab2:
    st.header("Bin Bay Mapping ↔️", divider='rainbow')
    st.markdown("Define bay definition groups and map bin IDs to bay types.")
    bay_types_list = [ "Bulk Stock", "Case Flow", "Drawer", "Flat Apparel", "Hanger Rod", "Hangers", "Jewelry", "Library", "Library Deep", "Pallet", "Shoes", "Random Other Bin", "PassThrough" ]
    bay_usage_options = [ "*", "45F Produce", "Aerosol", "Ambient", "Apparel", "BATTERIES", "BWS", "BWS_HIGH_FLAMMABLE", "BWS_LOW_FLAMMABLE", "BWS_MEDIUM_FLAMMABLE", "Book", "Chilled", "Chilled-FMP", "Corrosive", "Damage", "Damage Human Food", "Damage Pet Food", "Damage_HRV", "Damaged Aerosol", "Damaged Corrosive", "Damaged Flammable", "Misc Health Hazard", "Non Flammable Aerosols", "Non Inventory Storage-Facilities", "Non Inventory Storage-Other", "Non Inventory Storage-Stores", "Non Inventory-Black Totes", "Non Sort-Team Lift", "Non-Storage", "Non-TC Food", "Oxidizer", "Pet Food", "Produce", "Produce Backstock", "Produce Wetracks", "Reserve-Ambient", "Restricted Hazmat", "Semi-Chilled", "Shoes", "TC-Food", "Toxic", "Tropical" ]
    
    bay_groups_tab2 = []
    num_groups_tab2 = st.number_input("How many bay definition groups?", min_value=1, max_value=50, value=1, key="num_groups_bin_mapping")

    for group_idx in range(num_groups_tab2):
        if f"bin_group_name_{group_idx}" not in st.session_state:
            st.session_state[f"bin_group_name_{group_idx}"] = f"Bay Definition Group {group_idx + 1}"
        def update_bin_group_name(idx=group_idx):
            st.session_state[f"bin_group_name_{idx}"] = st.session_state[f"bin_group_name_input_{idx}"]
        header = st.session_state[f"bin_group_name_{group_idx}"].strip() or f"Bay Definition Group {group_idx + 1}"
        
        with st.expander(header, expanded=True):
            st.text_input("Group Name", value=st.session_state[f"bin_group_name_{group_idx}"], key=f"bin_group_name_input_{group_idx}", on_change=update_bin_group_name, args=(group_idx,))
            bin_ids_input = st.text_area(f"Enter bin IDs", key=f"bin_ids_{group_idx}", help="Paste Bin IDs from Excel (tab-separated, space-separated, or one per line).")
            bay_definition = st.text_input("Enter Bay Definition", max_chars=48, key=f"bay_definition_{group_idx}")
            st.divider()
            st.markdown("**Default Dimensions for the Group**")
            col1, col2, col3 = st.columns(3)
            with col1: height_cm = st.number_input("Height (CM)", min_value=0.0, value=0.0, key=f"height_cm_{group_idx}")
            with col2: width_cm = st.number_input("Width (CM)", min_value=0.0, value=0.0, key=f"width_cm_{group_idx}")
            with col3: depth_cm = st.number_input("Depth (CM)", min_value=0.0, value=0.0, key=f"depth_cm_{group_idx}")
            st.divider()
            outlier_shelves_input = st.text_input("Outlier Shelves (optional, comma-separated, e.g., C,D)", key=f"outlier_shelves_{group_idx}", help="Define shelves with different dimensions from the default.")
            st.caption("The app identifies a shelf by finding a capital letter followed by numbers at the end of the Bin ID (e.g., the 'C' in '...A208C120').")
            outlier_shelves = [s.strip().upper() for s in outlier_shelves_input.split(',') if s.strip()]
            outlier_dimensions = {}
            if outlier_shelves:
                for shelf in outlier_shelves:
                    st.markdown(f"**Dimensions for Outlier Shelf: {shelf}**")
                    o_col1, o_col2, o_col3 = st.columns(3)
                    with o_col1: o_height = st.number_input(f"Height (CM) for Shelf {shelf}", min_value=0.0, value=0.0, key=f"height_cm_{group_idx}_{shelf}")
                    with o_col2: o_width = st.number_input(f"Width (CM) for Shelf {shelf}", min_value=0.0, value=0.0, key=f"width_cm_{group_idx}_{shelf}")
                    with o_col3: o_depth = st.number_input(f"Depth (CM) for Shelf {shelf}", min_value=0.0, value=0.0, key=f"depth_cm_{group_idx}_{shelf}")
                    outlier_dimensions[shelf] = {"height_cm": o_height, "width_cm": o_width, "depth_cm": o_depth}
                st.divider()
            bay_usage = st.selectbox("Select Bay Usage", options=bay_usage_options, index=0, key=f"bay_usage_{group_idx}")
            bay_type = st.selectbox("Select Bay Type", options=bay_types_list, index=0, key=f"bay_type_{group_idx}")
            zone = st.text_input("Zone (e.g., Library (30D))", max_chars=25, key=f"zone_{group_idx}")
            if bin_ids_input:
                bin_list = [b.strip() for line in bin_ids_input.splitlines() for b in re.split(r'[\t\s]+', line) if b.strip()]
                if bin_list: bay_groups_tab2.append({"name": st.session_state[f"bin_group_name_{group_idx}"].strip(), "bin_ids": bin_list, "bay_definition": bay_definition, "height_cm": height_cm, "width_cm": width_cm, "depth_cm": depth_cm, "bay_usage": bay_usage, "bay_type": bay_type, "zone": zone, "outlier_dimensions": outlier_dimensions})
    
    if st.button("Generate Excel File", key="generate_bin_mapping_excel"):
        pass

with tab3:
    st.header("EOA Generator 🪧", divider='rainbow')
    st.markdown("**Step 1: Define All Aisles and Their Slot Ranges**")
    st.caption("Define all modules. For each, set a default slot range and specify any aisles with different slots.")
    num_mod_defs = st.number_input("How many modules do you want to define?", min_value=1, max_value=20, value=1, key="num_mod_defs")
    aisle_details = {} 
    for mod_idx in range(num_mod_defs):
        if f"eoa_mod_name_{mod_idx}" not in st.session_state:
            st.session_state[f"eoa_mod_name_{mod_idx}"] = ""
        def update_eoa_mod_name(idx=mod_idx):
            current_val = st.session_state[f"eoa_mod_name_input_{idx}"]
            st.session_state[f"eoa_mod_name_{idx}"] = current_val or f"Module Definition {idx + 1}"
        header = st.session_state[f"eoa_mod_name_{mod_idx}"] or f"Module Definition {mod_idx + 1}"
        with st.expander(header, expanded=True):
            mod_name = st.text_input("Module Name (e.g., P-1-A)", key=f"eoa_mod_name_input_{mod_idx}", on_change=update_eoa_mod_name, args=(mod_idx,)).strip()
            col1, col2 = st.columns(2)
            with col1: aisle_start = st.number_input(f"Start Aisle", min_value=1, value=200, step=1, key=f"aisle_start_{mod_idx}")
            with col2: aisle_end = st.number_input(f"End Aisle", min_value=aisle_start, value=aisle_start, step=1, key=f"aisle_end_{mod_idx}")
            st.divider()
            st.markdown("**Default Slot Range for this Module**")
            d_col1, d_col2 = st.columns(2)
            with d_col1: default_start_slot = st.number_input("Default Start Slot", value=1, step=1, key=f"d_slot_start_{mod_idx}")
            with d_col2: default_end_slot = st.number_input("Default End Slot", value=199, step=1, key=f"d_slot_end_{mod_idx}")
            outlier_aisles_input = st.text_area("Outlier Aisles for Slots (optional, comma-separated)", key=f"outlier_aisles_{mod_idx}")
            outlier_aisles = {int(a.strip()) for a in outlier_aisles_input.split(',') if a.strip()}
            outlier_slots = {}
            if outlier_aisles:
                st.markdown("**Outlier Slot Definitions**")
                for outlier in sorted(list(outlier_aisles)):
                    o_col1, o_col2 = st.columns(2)
                    with o_col1: outlier_start = st.number_input(f"Start Slot for Aisle {outlier}", value=1, step=1, key=f"o_start_{mod_idx}_{outlier}")
                    with o_col2: outlier_end = st.number_input(f"End Slot for Aisle {outlier}", value=199, step=1, key=f"o_end_{mod_idx}_{outlier}")
                    outlier_slots[outlier] = (outlier_start, outlier_end)
            if mod_name:
                aisle_details[mod_name] = {}
                aisles_in_range = list(range(aisle_start, aisle_end + 1))
                for aisle in aisles_in_range:
                    if aisle in outlier_slots:
                        aisle_details[mod_name][aisle] = {"slots": outlier_slots[aisle]}
                    else:
                        aisle_details[mod_name][aisle] = {"slots": (default_start_slot, default_end_slot)}
    st.divider()
    st.markdown("**Step 2: Define Physical Aisle Layouts**")
    st.markdown("**2a. Standard (Single-Module) Layouts**")
    st.caption("Describe how aisles within the same module are arranged.")
    standard_layout_input = st.text_area("Standard Layouts (one module per line)", height=150, key="eoa_standard_layout_input", placeholder="Example:\nP-1-A: 200, 201/202, 207")
    st.markdown("**2b. Cross-Module Pairs (Optional)**")
    st.caption("Define aisle pairs that touch across different modules.")
    cross_module_layout_input = st.text_area("Cross-Module Pairs (one pair per line)", height=100, key="eoa_cross_module_layout_input", placeholder="Example:\nP-1-A-201/P-1-B-200")
    st.divider()
    st.markdown("**Step 3: Confirm Placement Rule**")
    st.radio("Low End Placement Rule (for single-sided signs)", ["Odd on Left / Even on Right", "Even on Left / Odd on Right"], key="eoa_placement_rule", horizontal=True)
    if st.button("Generate EOA Signage", key="generate_eoa_signage"):
        # Generation Logic for Tab 3
        pass

with tab4:
    st.header("Bay Elevation Generator 📐", divider='rainbow')
    st.markdown("Define bay types, their shelves, and bin configurations to generate a PowerPoint elevation drawing.")
    st.info("ℹ️ **Note:** This feature requires `python-pptx`. Please add `python-pptx>=0.6.23` to your requirements.txt file.")
    num_bay_types = st.number_input("How many bay types do you want to define?", min_value=1, max_value=20, value=1, key="num_bay_types")
    bay_types_data = []
    for i in range(num_bay_types):
        if f"bay_type_name_{i}" not in st.session_state:
            st.session_state[f"bay_type_name_{i}"] = f"Bay Type {i + 1}"
        def update_bay_type_name(idx=i): st.session_state[f"bay_type_name_{idx}"] = st.session_state[f"bay_type_name_input_{idx}"]
        header = st.session_state[f"bay_type_name_{i}"].strip() or f"Bay Type {i + 1}"
        with st.expander(header, expanded=True):
            bay_name = st.text_input("Bay Type Name", value=st.session_state[f"bay_type_name_{i}"], key=f"bay_type_name_input_{i}", on_change=update_bay_type_name, args=(i,))
            shelf_count = st.number_input("Number of shelves in this bay?", min_value=1, max_value=26, value=3, key=f"elevation_shelf_count_{i}")
            shelf_names = list(string.ascii_uppercase[:shelf_count])
            shelf_details = {}
            for shelf_name in shelf_names:
                st.divider()
                st.markdown(f"**Configuration for Shelf {shelf_name}**")
                num_bins = st.number_input("Number of bins in this shelf?", min_value=1, max_value=50, value=5, key=f"num_bins_{i}_{shelf_name}")
                st.markdown(f"**Bin Dimensions for all bins in Shelf {shelf_name} (cm)**")
                c1, c2, c3 = st.columns(3)
                with c1: bin_h = st.number_input("Height", min_value=1.0, value=10.0, key=f"bin_h_{i}_{shelf_name}")
                with c2: bin_w = st.number_input("Width", min_value=1.0, value=10.0, key=f"bin_w_{i}_{shelf_name}")
                with c3: bin_d = st.number_input("Depth", min_value=1.0, value=10.0, key=f"bin_d_{i}_{shelf_name}")
                shelf_details[shelf_name] = {'num_bins': num_bins, 'h': bin_h, 'w': bin_w, 'd': bin_d}
            if bay_name.strip():
                bay_types_data.append({"name": bay_name.strip(), "shelves": shelf_names, "shelf_details": shelf_details})
    if st.button("Generate PowerPoint Elevation", key="generate_ppt"):
        if any(not bay['name'] for bay in bay_types_data):
            st.error("Please provide a name for all bay types.")
        else:
            with st.spinner("Generating PowerPoint file..."):
                try:
                    ppt_buffer = generate_elevation_powerpoint(bay_types_data)
                    st.success(f"✅ Success! Generated a PowerPoint with {len(bay_types_data)} slides.")
                    st.download_button(
                        label="📥 Download PowerPoint File",
                        data=ppt_buffer,
                        file_name="bay_elevations.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
                except Exception as e:
                    st.error(f"An error occurred during PowerPoint generation: {e}")
